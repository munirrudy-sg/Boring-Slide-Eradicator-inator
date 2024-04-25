from pptx import Presentation

prs = Presentation("Theme1.pptx")
slide = prs.slides.add_slide(prs.slide_layouts[0])
for shape in slide.placeholders:
  print('%d %s' % (shape.placeholder_format.idx, shape.name))