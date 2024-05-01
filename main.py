import time
import os
import streamlit as st
from dotenv import load_dotenv
import google.generativeai as gen_ai
import fitz  # pip install PyMuPDF
from pptx import Presentation  # pip install python-pptx
from pptx.util import Pt, Inches
import io
import re
from streamlit_image_select import image_select


def convert_slides_data_to_text(slides_data):
    # Create a plain text representation of slides_data
    slides_text = ""
    for slide in slides_data:
        slide_number, title, content = slide
        slides_text += f"Slide {slide_number}:\n"
        slides_text += f"Title: {title}\n"
        slides_text += "Content:\n"
        slides_text += content + "\n"  # Adding a new line for separation
        slides_text += "\n" + "-" * 40 + "\n"  # Separator between slides
    return slides_text

# Function to open the PDF and store it in session state
def store_pdf_in_session(uploaded_file):
    # Check if a file has been uploaded
    if uploaded_file is not None:
        try:
            # Ensure the file pointer is at the start
            uploaded_file.seek(0)
            # Open the PDF and store it in session state
            doc = fitz.open("pdf", uploaded_file.read())
            st.session_state['pdf_doc'] = doc
        except fitz.EmptyFileError:
            st.error("Cannot open the PDF stream. Please upload a valid PDF.")
    else:
        # # Handle the case where no file is uploaded
        # Reset or clear any existing PDF data in the session state
        if 'pdf_doc' in st.session_state:
            del st.session_state['pdf_doc']



def apply_bold_format(word, r, bolded_mode):
    # Check if the word starts and ends with double asterisks
    if word.startswith("**") and word.endswith("**"):
        r.text = word[2:-2]  # Remove outer asterisks
        r.font.bold = True
    elif word.startswith("**"):
        r.text = word[2:] + " "  # Remove leading asterisks
        r.font.bold = True
        bolded_mode = True  # Switch to bolded mode
    elif word.endswith("**"):
        r.text = word[:-2] + " "  # Remove trailing asterisks
        r.font.bold = True
        bolded_mode = False  # Exit bolded mode
    elif bolded_mode:
        r.text = word + " "  # Apply bolded mode to other words
        r.font.bold = True
    else:
        r.text = word + " "  # Normal text, no bold
    return bolded_mode


def create_slide(prs, title, content):
    # Retrieve the theme selected by the user
    selected_theme = st.session_state.get("selected_theme", "Theme1")
    layout, text_placeholder = theme_dict.get(selected_theme, (11, 14))  # Default to (11, 14)

    # Create a new slide with appropriate layout
    slide_layout = prs.slide_layouts[layout]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)

    # Set the title text box
    title_shape = slide.shapes.title
    title_shape.text = title  # Assign the title

    # Use the content placeholder to set the content text
    content_shape = slide.placeholders[text_placeholder]
    content_text_frame = content_shape.text_frame
    content_text_frame.word_wrap = True  # Ensure proper word wrapping

    # Split content into lines based on newline character
    lines = content.replace("\\n", "\n").split("\n") # Convert escaped newlines


    for line in lines:
        if line.startswith("##"):
            # Heading line
            p = content_text_frame.add_paragraph()
            p.text = line[2:].strip()  # Remove "##"
            p.font.size = Pt(20)  # Larger font for headings
            p.font.bold = True
        elif line.startswith("* "):
            # Bullet point
            p = content_text_frame.add_paragraph()
            line=line[1:].strip()
            words = line.split()
            bolded_mode = False
            for word in words:
                r = p.add_run()
                bolded_mode = apply_bold_format(word, r, bolded_mode)  # Apply bolding
            for run in p.runs:
                run.font.size = Pt(12)  # Smaller font for bullet points
        else:
            # Normal text
            p = content_text_frame.add_paragraph()
            words = line.split()
            bolded_mode = False
            for word in words:
                r = p.add_run()
                bolded_mode = apply_bold_format(word, r, bolded_mode)  # Apply bold if needed
                r.font.size = Pt(14)  # Standard font size for regular text


def create_presentation(data):
    # Retrieve the theme selected by the user
    selected_theme = st.session_state.get("selected_theme", "Theme1")

    # Create the absolute path to the theme
    theme_dir = os.path.abspath("theme_pptx")
    theme_path = os.path.join(theme_dir, f"{selected_theme}.pptx")
    
    # Append ".pptx" to create the theme path
    # theme_path = f"theme_pptx\\{selected_theme}.pptx"
    prs = Presentation(theme_path)

    for slide_data in data:
        title = slide_data.get("title", "")
        content = slide_data.get("content", "")
        create_slide(prs, title, content)

    return prs

# Dictionary to store theme-based layout and text_placeholder mappings
theme_dict = {
    "Theme1": (0, 1),
    "Theme2": (11, 14),
    "Theme3": (0, 10),
    "Theme4": (0, 10),
}

def process_pdf(uploaded_file):
    # Process PDF code (your existing code)
    status_message = st.empty()
    my_progress = st.progress(0, text="Activating the Boring-Slide-Eradicator-inator...")

    doc = fitz.open("pdf", uploaded_file.read())  # Load the PDF file
    slides_data = []

# List to store titles extracted from PDF
    titles = []
    original_content = []

    # Process each page
    total_pages = len(doc)
    for i, page in enumerate(doc):
        text = page.get_text()
        lines = text.split('\n')
        title = lines[0] if lines else "No Title"
        content = "\n".join(lines[1:]) if len(lines) > 1 else "No Content"

        slides_data.append((i + 1, title, content))
        titles.append(title)
        original_content.append(content)

        progress = (i + 1) / total_pages * 25
        my_progress.progress(int(progress), "Perry the Platypus is reviewing your slides!")
        time.sleep(1)


    st.session_state.slides_data = slides_data

    # Exclude selected slides

    # Get the index of the selected options in slide_options
    exclude_slide_indices = [slide_options.index(option) for option in exclude_slides]
    exclude_slide_numbers = [index + 1 for index in exclude_slide_indices]

    st.session_state.exclude_slide_numbers = exclude_slide_numbers

    # Convert slides_data to plain text
    slides_text = convert_slides_data_to_text(slides_data)

    # Construct the prompt for the AI model
    if custom_prompt_text:
        prompt = f"The below text is extracted text from lecture slides. Based on the input text '{custom_prompt_text}', prioritize knowledge accordingly. Currently, the slides are bad and I want you to replace each slide content with your more well-explained version. Don't just explain it point by point. I want you to understand what the slide is trying to say then explain it in a way that can be easily understood by a university student. Include citations. Don't just write in plain text, use bullet points or anything that makes the student read and understand the slide easily. Do this for every slide but don't add new slides. Your output should be the slide number, title and slide content. Each slide separated by comma and text in markdown.\n\n.{slides_text}"
    else:
        prompt = f"The below text is extracted text from lecture slides. Currently, the slides are bad and I want you to replace each slide content with your more well-explained version. Don't just explain it point by point. I want you to understand what the slide is trying to say then explain it in a way that can be easily understood by a university student. Include citations. Don't just write in plain text, use bullet points or anything that makes the student read and understand the slide easily. Do this for every slide but don't add new slides. Your output should be the slide number, title and slide content. Each slide separated by comma and text in markdown.\n\n.{slides_text}"

    # Loop to manage retries
    retry_count = 0
    max_retries = 3
    success = False

    while retry_count < max_retries and not success:

        for i in range(25, 80):
            if retry_count == 0:
                my_progress.progress(i, "Supercharging your presentation with Gemini juice!")
                time.sleep(1)
            else:
                my_progress.progress(i, "Uh oh! Resending prompt due to unexpected error...")
                time.sleep(1)

        # Send the prompt to the AI model
        response = model.generate_content(prompt, request_options={"timeout": 600000})

        content_pattern = r'"content":\s*"([^"]+)"'
        content_matches = re.findall(content_pattern, response.text, re.DOTALL)

        # Create a new list with PDF slide titles and AI response slide content
        PPT_data = []
        if len(titles) == len(content_matches):
            for i, content in enumerate(content_matches):
                slide_number = i + 1
                title = titles[i]
                original_content_text = original_content[i]
                if slide_number in exclude_slide_numbers:
                    PPT_data.append({
                        "slide_number": slide_number,
                        "title": title,
                        "content": original_content_text
                    })
                else:
                    PPT_data.append({
                        "slide_number": slide_number,
                        "title": title,
                        "content": content
                    })
            
            # Create PowerPoint presentation and save it
            presentation = create_presentation(PPT_data)
            output_data = io.BytesIO()
            presentation.save(output_data)
            output_data.seek(0)
            st.session_state.processed = True
            st.session_state.output_data = output_data
            st.session_state.output_data = open("output.pptx", "rb").read() # to delete

            # Clear the status message once done processing
            status_message.empty()
            my_progress.progress(100, "And POOF! Your slides is slides-tastic!")

            success = True

            # Return the output data
            return output_data

        else:
            # Indicate mismatch and retry
            retry_count += 1

    if retry_count == max_retries and not success:
        status_message.empty()
        my_progress.progress(100, "Error occured! This usually never happens but let's try again in a moment.")
        return None

# Load environment variables
load_dotenv()

# Configure Streamlit page settings
st.set_page_config(
    page_title="Boring-Slide-Eradicator-inator",
    page_icon="🧙‍♂️",  # Favicon emoji
    layout="centered",  # Page layout option
)

GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY")

# Set up Google Gemini-Pro AI model
gen_ai.configure(api_key=GOOGLE_API_KEY)

# Set up the model
generation_config = {
    "temperature": 1,
    "top_p": 0.95,
    "top_k": 0,
    "max_output_tokens": 200000,
    "response_mime_type": "application/json",
}

safety_settings = [
    {
        "category": "HARM_CATEGORY_HARASSMENT",
        "threshold": "BLOCK_NONE"
    },
    {
        "category": "HARM_CATEGORY_HATE_SPEECH",
        "threshold": "BLOCK_NONE"
    },
    {
        "category": "HARM_CATEGORY_SEXUALLY_EXPLICIT",
        "threshold": "BLOCK_NONE"
    },
    {
        "category": "HARM_CATEGORY_DANGEROUS_CONTENT",
        "threshold": "BLOCK_NONE"
    },
]

model = gen_ai.GenerativeModel(model_name="gemini-1.5-pro-latest",
                               generation_config=generation_config,
                               safety_settings=safety_settings)

# Display the app's title on the page
st.title("🧙‍♂️ Boring-Slide-Eradicator-inator")

# Sidebar
st.sidebar.markdown("""
## Greetings! 👋

Is staring at confusing lecture slides making you feel like Perry the Platypus trapped in a cage? ‍  Been there, done that!

That's why we invented the **Boring-Slide-Eradicator-inator**!   This super-cool Gemini-powered tool uses some serious science to transform those cryptic slides into clear, understandable slides. ✨

### Here's how this inator works its magic:
- **Import your PDF files**: Toss in those troublesome lecture slides.
- **Give Gemini some context**: Tell it what materials to reference for explanations.
- **Exclude unnecessary slides**: Skip the title pages and "thank you" slides.
- **Choose your theme**: Because even explanations deserve a little style.
- **Hit "Process PDF"** and watch the magic happen! 

Get ready to say goodbye to slide-induced headaches and hello to understanding!
""")


# Check if the session state already has a 'processed' flag
if 'processed' not in st.session_state:
    st.session_state.processed = False

if 'slides_data' not in st.session_state:
    st.session_state.slides_data = []
uploaded_file = st.file_uploader("Step 1: Choose a boring PDF file", type=['pdf'])
store_pdf_in_session(uploaded_file)

# Text input for custom prompt
custom_prompt_text = st.text_input("Step 2: Enter materials to refer from (optional):", help="Enter your textbooks here if you want Gemini to use them")

# Ensure the file has been uploaded and is not empty
if uploaded_file is None or uploaded_file.size == 0:  # If file is removed or empty
    st.session_state.processed = False  # Reset the processed flag
    st.session_state.slides_data = []  # Reset slides data if needed
    st.session_state.output_data = None  # Reset output data
    st.session_state.exclude_slide_numbers = []
    store_pdf_in_session(uploaded_file)

else:  # File has been uploaded
    if not st.session_state.processed:
        store_pdf_in_session(uploaded_file)

# Use the stored PDF to populate the dropdown for excluding slides
if 'pdf_doc' in st.session_state:
    doc = st.session_state['pdf_doc']
    # Create slide options with slide number and title only
    slide_options = []
    for i in range(len(doc)):
        title = doc[i].get_text("text").split("\n")[0]  # First line as title
        slide_options.append(f"Slide {i + 1}: {title}")
else:
    slide_options = []

exclude_slides = st.multiselect("Step 3: Select slides to exclude (optional):", slide_options, help="These slides would not be enhanced by Gemini")

selected_theme = image_select(
    label="Step 4: Spice up your PPT",
    images=[f"theme_thumbnails\\{theme}.jpg" for theme in theme_dict]
)

# Create a button to process the PDF
process_button = st.button("Zap that PDF!")

# When the button is pressed
if process_button:
    if uploaded_file is None:
        st.error("Please upload a PDF before processing.")
    else:
        # Reset the file pointer to the beginning
        uploaded_file.seek(0)
        # Save the selected theme in session state
        st.session_state.selected_theme = str(selected_theme.split("\\")[-1].split(".")[0])

        # Process the PDF if the file is uploaded
        output_data = process_pdf(uploaded_file)
        st.session_state.processed = True

    # Provide a download button for the processed output
    if 'output_data' in st.session_state:
        st.success('Processing complete!')
        st.download_button(label="Download PowerPoint", data=st.session_state.output_data, file_name="output.pptx",
                           mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")