import json
from pptx import Presentation
from pptx.util import Pt,Inches

# text = '''- **Reverse Engineering in the Real World:**\\n   * Historically, software reverse engineering has been at the center of numerous legal disputes. Companies like Atari, Accolade, Phoenix, and Connectix engaged in reverse engineering to
# #  achieve compatibility with competitor's systems, leading to landmark cases.  \\n   * These cases often revolved around copyright infringement and the extent to which companies could analyze and utilize competitors' software to create compatible products. \\n- **Legal Precedents:**\\n   * Interestingly, man
# # y of these cases were ruled in favor of reverse engineering, establishing a legal basis for the practice under certain conditions.  \\n   * Courts recognized the importance of interoperability and innovation, allowing companies to reverse engineer software for legitimate purposes such as compatibility, secu
# # rity analysis, and research. \\n- **Challenges and the Chinese Wall Method:**\\n   * Proving that reverse engineering was done without infringing copyright can be challenging.  \\n   * The 'Chinese Wall' method emerged as a strategy to address this concern. It involves separating the team performing reverse engineering from the development team to prevent the direct use of copyrighted material. \\n   * This method aimed to demonstrate good faith efforts to avoid copyright infringement during the reverse engineering process.'''


def apply_bold_format(word, r, bolded_mode):
    # Check if the word starts and ends with double asterisks
    if word.startswith("**") and word.endswith("**"):
        r.text = word[2:-2]  # Remove outer asterisks
        r.font.bold = True
    elif word.startswith("**"):
        r.text = word[2:] + " "  # Remove leading asterisks
        r.font.bold = True
        bolded_mode = True  # Switch to bolded mode
    elif word.endswith("**"):
        r.text = word[:-2] + " "  # Remove trailing asterisks
        r.font.bold = True
        bolded_mode = False  # Exit bolded mode
    elif bolded_mode:
        r.text = word + " "  # Apply bolded mode to other words
        r.font.bold = True
    else:
        r.text = word + " "  # Normal text, no bold
    return bolded_mode


def create_slide(prs, title, content):
    # Create a new slide with appropriate layout
    slide_layout = prs.slide_layouts[5]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)

    # Set the title text box
    title_shape = slide.shapes.title
    title_shape.text = title  # Assign the title

    # Create the content text box
    content_shape = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
    content_text_frame = content_shape.text_frame
    content_text_frame.word_wrap = True  # Ensure proper word wrapping

    # Split content into lines based on newline character
    lines = content.replace("\\n", "\n").split("\n")  # Convert escaped newlines

    for line in lines:
        if line.startswith("##"):
            # Heading line
            p = content_text_frame.add_paragraph()
            p.text = line[2:].strip()  # Remove "##"
            p.font.size = Pt(20)  # Larger font for headings
            p.font.bold = True
        elif line.startswith("*"):
            # Bullet point
            p = content_text_frame.add_paragraph()
            words = line.split()
            bolded_mode = False
            for word in words:
                r = p.add_run()
                bolded_mode = apply_bold_format(word, r, bolded_mode)  # Apply bolding
            for run in p.runs:
                run.font.size = Pt(12)  # Smaller font for bullet points
        else:
            # Normal text
            p = content_text_frame.add_paragraph()
            words = line.split()
            bolded_mode = False
            for word in words:
                if(word=="**Challenges"):
                    print("ASD")
                r = p.add_run()
                apply_bold_format(word, r, bolded_mode)  # Apply bold if needed
                r.font.size = Pt(14)  # Standard font size for regular text


def create_presentation(data):
    prs = Presentation()

    # Set slide width and height for 16:9 aspect ratio
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    for slide_data in data:
        title = slide_data.get("title", "")
        content = slide_data.get("content", "")
        create_slide(prs, title, content)

    return prs

# Example usage:
data=[
    {
        "slide_number": 1,
        "title": "Android Reverse Engineering",
        "content": '''
        - **Reverse Engineering in the Real World:**\\n   * Historically, software reverse engineering has been at the center of numerous legal disputes. Companies like Atari, Accolade, Phoenix, and Connectix engaged in reverse engineering to
#  achieve compatibility with competitor's systems, leading to landmark cases.  \\n   * These cases often revolved around copyright infringement and the extent to which companies could analyze and utilize competitors' software to create compatible products. \\n- **Legal Precedents:**\\n   * Interestingly, man
# y of these cases were ruled in favor of reverse engineering, establishing a legal basis for the practice under certain conditions.  \\n   * Courts recognized the importance of interoperability and innovation, allowing companies to reverse engineer software for legitimate purposes such as compatibility, secu
# rity analysis, and research. \\n- **Challenges and the Chinese Wall Method:**\\n   * Proving that reverse engineering was done without infringing copyright can be challenging.  \\n   * The 'Chinese Wall' method emerged as a strategy to address this concern. It involves separating the team performing reverse engineering from the development team to prevent the direct use of copyrighted material. \\n   * This method aimed to demonstrate good faith efforts to avoid copyright infringement during the reverse engineering process.
        '''
    },
    {
        "slide_number": 2,
        "title": "Overview",
        "content": "## Course Roadmap\n\n* **Android Reverse Engineering:** Understanding the core principles and methodologies.\n* **Android Code Structure:** Exploring the organization of an Android app's codebase.\n* **Components in Manifest file:** Decoding the essential app metadata within the manifest.\n* **Dex to Smali:** Disassembling DEX bytecode into human-readable Smali code.\n* **Dex to Source code:** Converting DEX bytecode back into Java or Kotlin source code.\n* **Native code analysis:** Examining native code components written in languages like C.\n* **Design a dummy machine language:** A practical exercise to solidify your understanding of low-level code and assembly language."
    },
    {
        "slide_number": 3,
        "title": "Administrative Matters",
        "content": "## Mark Your Calendars!\n\n* **Guest Lecture:** We'll have a special session led by DarkNavy, a renowned security expert from China.\n* **Date & Time:** 19th March 2024, 9 AM – 11 AM\n* **Format:** Hybrid - Attend in person or join via Zoom (classroom details will be on LMS).\n* **CTF Session:** Get hands-on experience with a Capture the Flag challenge to hone your skills and identify your strengths."
    },
    {
        "slide_number": 4,
        "title": "Introduction to Android Reverse Engineering",
        "content": "## Demystifying Android Apps\n\nReverse engineering is like taking apart a complex machine to understand its inner workings. In the context of Android apps, it involves analyzing the compiled code and resources to:\n\n* **Comprehend App Behavior:** Discover how the app functions at a deeper level.\n* **Identify Security Flaws:** Uncover vulnerabilities that could be exploited by malicious actors.\n* **Customize and Modify:**  Potentially create modifications or enhancements to the app (always respecting ethical and legal boundaries)."
    },
    {
        "slide_number": 5,
        "title": "Android App Structure",
        "content": "## Inside the APK\n\nAndroid apps are packaged as APK (Android Package) files, which are essentially zip archives containing everything needed to run the app. Key components include:\n\n* **AndroidManifest.xml:** Defines the app's structure, components, permissions, and other essential information.\n* **res/ (Resources):** Stores images, layouts, strings, and other resources used by the app.\n* **assets/ (Assets):** Contains raw asset files, such as fonts or game data.\n* **DEX files:** Holds the compiled Java/Kotlin code (Dalvik Executable format)."
    },
    {
        "slide_number": 6,
        "title": "Android Manifest",
        "content": "## The App Blueprint\n\nThe AndroidManifest.xml file serves as the app's blueprint, declaring its core components and capabilities:\n\n* **Activities:** The building blocks of the user interface, representing individual screens.\n* **Services:** Background processes that perform long-running operations.\n* **Content Providers:** Manage shared data and facilitate data access between apps.\n* **Broadcast Receivers:** Respond to system-wide events or broadcasts.\n* **Intent Filters:** Define how activities, services, and receivers handle specific intents.\n* **Permissions:** Specify the app's access to system resources and user data."
    },
    {
        "slide_number": 7,
        "title": "Android Manifest components",
        "content": "## Manifest Essentials\n\nHere's a breakdown of key elements within the manifest:\n\n* **manifest:** The root element, defining the package name and namespace declarations.\n* **uses-sdk:** Specifies the minimum and target SDK versions supported by the app.\n* **uses-permission:** Declares the permissions required by the app, such as access to the camera or internet."
    },
    {
        "slide_number": 8,
        "title": "Android Manifest components",
        "content": "## Application Attributes\n\nThe application element houses app-level settings:\n\n* **android:name:** Specifies the application class name.\n* **android:allowBackup:** Controls whether backups are allowed.\n* **android:icon and android:label:** Define the app's icon and name.\n* **android:debuggable:** Enables debugging during development."
    },
    {
        "slide_number": 9,
        "title": "Android Manifest components",
        "content": "## Activities and Intents\n\n* **activity:** Declares an activity within the app. Activities must be declared to be functional.\n* **intent-filter:** Specifies the type of intents the component can respond to. The `MAIN` action and `LAUNCHER` category designate the app's entry point."
    },
    {
        "slide_number": 10,
        "title": "What to look for?",
        "content": "## Exported Activities: A Potential Entry Point\n\nPay close attention to the `android:exported` attribute within activity declarations. If set to `true`, the activity can be launched from outside the app, potentially exposing sensitive functionality. Use tools like `adb shell` to explore and interact with such activities."
    },
    {
        "slide_number": 11,
        "title": "Reverse Engineering",
        "content": "## Decoding the Code\n\nAndroid apps typically use Java or Kotlin, which are compiled into DEX (Dalvik Executable) bytecode. Reverse engineering involves:\n\n* **Dex to Smali:** Converting DEX into Smali, a human-readable intermediate language, using tools like apktool.\n* **Native Code Analysis:** If the app contains native code (C/C++), it needs to be analyzed at the assembly level using disassemblers and debuggers."
    },
    {
        "slide_number": 12,
        "title": "Native Code Analysis",
        "content": "## Bridging the Gap\n\nNative code is often found in libraries within the app's resources directory. To understand its interaction with the Java/Kotlin code, look for methods like `loadLibrary` or `load` that indicate native function calls. While the native code itself requires separate analysis, identifying these calls provides valuable insights."
    },
    {
        "slide_number": 13,
        "title": "Understanding code logic (Dex to Java)",
        "content": "## Recreating the Source\n\nTools like JADX, dex2jar, and online services like https://decompiler.com can attempt to convert DEX bytecode back into Java source code. However, the process isn't perfect and may not always produce compilable code. Despite limitations, decompilation can provide valuable insights into the app's logic and algorithms.\n\nFor Kotlin-based apps, consider using Fernflower for decompilation."
    },
    {
        "slide_number": 14,
        "title": "Modifying and Patching App",
        "content": "## App Modification: Tread Carefully\n\nModifying an app involves several steps:\n\n1. **Reverse to Smali:** Convert the DEX code into Smali format.\n2. **Edit Smali Code:** Make the desired changes directly within the Smali files.\n3. **Modify Resources (Optional):** Alter resources like images or layouts as needed.\n4. **Recompile and Re-sign:** Rebuild the APK and sign it with a valid certificate.\n\n**Important:** Always respect the app's End User License Agreement (EULA) and copyright laws. Modifying apps without permission can have legal consequences."
    },
    {
        "slide_number": 15,
        "title": "Crash Course on Programming Language",
        "content": "## Building Blocks of Code\n\nTo solidify your understanding of low-level code and assembly language, we'll embark on a practical journey:\n\n* **Designing a Dummy Assembly Language:** We'll create a simplified assembly language with its own set of instructions.\n* **Binary Code Representation:** We'll define how these instructions are represented in binary form.\n* **Assembly to Binary Conversion:** Learn how to convert assembly code into its corresponding binary representation."
    }
]
presentation = create_presentation(data)
presentation.save("presentation.pptx")